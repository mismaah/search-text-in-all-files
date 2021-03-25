import os
import docx2txt
import zipfile
import re
import time
from pptx import Presentation
from pdfminer import high_level

CONTEXT_RANGE = 3


def display(text, results, file):
    plural = "es" if len(results) > 1 else ""
    print(f"{len(results)} match{plural} in {file}:")
    for i in results:
        context = text[i[0] - 20 : i[1] + 20].strip()
        context = context.replace("\n", " ")
        print(f"\t{context}")
    print("")


supportedFormats = ["docx", "txt", "pptx", "pdf"]
filePaths = []
for root, _, files in os.walk(".\\"):
    for f in files:
        if f.split(".")[-1] in supportedFormats:
            filePaths.append([f"{root}\{f}", f.split(".")[-1]])
while True:
    matches = 0
    query = input("Search for: ").lower()
    start = time.time()
    for i in filePaths:
        text = ""
        if i[1] == "docx":
            try:
                text = docx2txt.process(i[0])
            except zipfile.BadZipFile:
                print(f"Cannot read {i[0]}.")
        elif i[1] == "pptx":
            prs = Presentation(i[0])
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        elif i[1] == "pdf":
            text = high_level.extract_text(i[0])
        else:
            with open(i[0], "r", encoding="utf-8") as f:
                try:
                    text = f.read()
                except UnicodeDecodeError:
                    print(f"Cannot read {i[0]}.")
        file_results = []
        for match in re.finditer(query, text.lower()):
            matches += 1
            file_results.append(match.span())
        if not file_results:
            continue
        display(text, file_results, i[0])
    if matches == 0:
        print("No results.")
    print(f"Time elapsed for search: {int((time.time() - start) * 1000)}ms")
