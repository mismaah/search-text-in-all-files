import os
import docx2txt
import zipfile
import re
import time
from pptx import Presentation
from pdfminer import high_level

CONTEXT_RANGE = 3
MATCH_OVER_SIMILARITY = 0.65


def generate_trigrams(text: str):
    text = text.lower()
    text = text.strip()
    text = re.sub(" ", "  ", text)
    text = f"  {text} "
    return set([text[i : i + 3] for i in range(len(text) - 2)])


def compare_trigrams(trig1: set, trig2: set):
    count = 0
    for i in trig1:
        if i in trig2:
            count += 1
    return count / len(trig2)


def trigram_search(query, file_text, file):
    query_trig = generate_trigrams(query)
    results = []
    tokens = file_text.split()
    for j in range(len(tokens)):
        word_trig = generate_trigrams(tokens[j])
        similarity = compare_trigrams(query_trig, word_trig)
        if similarity < MATCH_OVER_SIMILARITY:
            continue
        before = ""
        after = ""
        for k in range(1, CONTEXT_RANGE):
            if j - k > 0:
                before = f"{tokens[j-k]} " + before
            if j + k < len(tokens):
                after += f" {tokens[j+k]}"
        result = {
            "word": tokens[j],
            "similarity": similarity,
            "context": f"{before}{tokens[j]}{after}",
        }
        results.append(result)
    if not results:
        return
    results_combined = []
    for j in set([(i["word"], i["similarity"]) for i in results]):
        contexts = []
        for result in results:
            if j[0] == result["word"] and j[1] == result["similarity"]:
                contexts.append(result["context"])
        combined = {
            "file": file,
            "word": j[0],
            "similarity": j[1],
            "contexts": contexts,
        }
        results_combined.append(combined)
    return results_combined


def display(results):
    results = sorted(results, key=lambda k: k["similarity"], reverse=True)
    for i in results:
        print(
            f"{int(round(i['similarity'] * 100))}% match for {i['word']} in [{i['file'][0]}]:"
        )
        for j in i["contexts"]:
            print(f"\t{j}")
        print("")


supportedFormats = ["docx", "txt", "pptx", "pdf"]
filePaths = []
for root, _, files in os.walk(".\\"):
    for f in files:
        if f.split(".")[-1] in supportedFormats:
            filePaths.append([f"{root}\{f}", f.split(".")[-1]])
while True:
    matches = 0
    results = []
    query = input("Search for: ").lower()
    start = time.time()
    for i in filePaths:
        print(f"Searching {i[0]}...")
        text = ""
        if i[1] == "docx":
            try:
                text = docx2txt.process(i[0])
            except zipfile.BadZipFile:
                print(f"Cannot read {i[0]}.")
        elif i[1] == "pptx":
            prs = Presentation(i[0])
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        elif i[1] == "pdf":
            text = high_level.extract_text(i[0])
        else:
            with open(i[0], "r", encoding="utf-8") as f:
                try:
                    text = f.read()
                except UnicodeDecodeError:
                    print(f"Cannot read {i[0]}.")
        file_results = trigram_search(query, text, i)
        if not file_results:
            continue
        results += file_results
    if len(results) == 0:
        print("No results.")
    else:
        display(results)
    print(f"Time elapsed for search: {int((time.time() - start) * 1000)}ms")
