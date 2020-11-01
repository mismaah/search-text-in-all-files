import os
import docx2txt
import zipfile
from pptx import Presentation
supportedFormats = ['docx', 'txt', 'pptx']
filePaths = []
for root, _, files in os.walk('.\\'):
    for f in files:
        if f.split('.')[-1] in supportedFormats:
            filePaths.append([f'{root}\{f}', f.split('.')[-1]])
while True:
    query = input("Search for: ").lower()
    results = []
    for i in filePaths:
        matches = 0
        contexts = []
        text = ""
        if i[1] == 'docx':
            try:
                text = docx2txt.process(i[0])
            except zipfile.BadZipFile:
                print(f"Cannot read {i[0]}.")
        elif i[1] == 'pptx':
            prs = Presentation(i[0])
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        else:
            with open(i[0], "r", encoding="utf-8") as f:
                try:
                    text = f.read()
                except UnicodeDecodeError:
                    print(f"Cannot read {i[0]}.")
        tokens = text.split()
        for j in range(len(tokens)):
            if query in tokens[j].lower():
                matches += 1
                if tokens[j] not in contexts:
                    contextRange = 3
                    before = ""
                    after = ""
                    for k in range(1, contextRange):
                        if j-k > 0:
                            before = f"{tokens[j-k]} " + before
                        if j+k < len(tokens):
                            after += f" {tokens[j+k]}"
                    contexts.append(f"{before}{tokens[j]}{after}")
        if matches > 0:
            result = {
                "file": i,
                "matches": matches,
                "contexts": contexts
            }
            results.append(result)
    if len(results) == 0:
        print("No results.")
    results = sorted(results, key=lambda k: k["matches"], reverse=True)
    for i in results:
        plural = "" if i['matches'] == 1 else "es"
        print(f"{i['matches']} match{plural} in [{i['file'][0]}]")
        for j in i["contexts"]:
            print(f"\t{j}")
        print("")